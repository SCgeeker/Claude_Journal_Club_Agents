"""
投影片生成器模組
基於Journal Club的多風格學術簡報生成
"""

import re
import requests
import json
import os
from pathlib import Path
from typing import Dict, List, Optional, Any, Tuple
from datetime import datetime
import time
import yaml as yaml_lib

try:
    from jinja2 import Template
    JINJA2_AVAILABLE = True
except ImportError:
    JINJA2_AVAILABLE = False

try:
    import yaml
    YAML_AVAILABLE = True
except ImportError:
    YAML_AVAILABLE = False

# 嘗試導入模型監控模塊
try:
    from src.utils.model_monitor import ModelMonitor
    MONITOR_AVAILABLE = True
except ImportError:
    try:
        # 備用導入路徑
        import sys
        sys.path.insert(0, str(Path(__file__).parent.parent.parent))
        from src.utils.model_monitor import ModelMonitor
        MONITOR_AVAILABLE = True
    except ImportError:
        MONITOR_AVAILABLE = False
        ModelMonitor = None

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# 可選的LLM後端
try:
    import openai
    OPENAI_AVAILABLE = True
except ImportError:
    OPENAI_AVAILABLE = False

try:
    import google.generativeai as genai
    GOOGLE_AVAILABLE = True
except ImportError:
    GOOGLE_AVAILABLE = False

try:
    import anthropic
    ANTHROPIC_AVAILABLE = True
except ImportError:
    ANTHROPIC_AVAILABLE = False


class SlideMaker:
    """
    投影片生成器
    支援7種學術風格、5種詳細程度、3種語言
    """

    def __init__(self,
                 template_path: Optional[str] = None,
                 styles_config: Optional[str] = None,
                 llm_provider: str = "auto",
                 ollama_url: str = "http://localhost:11434",
                 api_key: Optional[str] = None,
                 selection_strategy: str = "balanced"):
        """
        初始化投影片生成器

        Args:
            template_path: Jinja2模板路徑
            styles_config: 風格配置YAML路徑
            llm_provider: LLM提供者 (auto/ollama/openai/google/anthropic)
            ollama_url: Ollama API地址
            api_key: API金鑰（OpenAI/Google/Anthropic用）
            selection_strategy: 模型選擇策略 (balanced/quality_first/cost_first/speed_first)
        """
        if not JINJA2_AVAILABLE:
            raise ImportError("Jinja2 not installed. Run: pip install jinja2")

        if not YAML_AVAILABLE:
            raise ImportError("PyYAML not installed. Run: pip install pyyaml")

        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx not installed. Run: pip install python-pptx")

        self.llm_provider = llm_provider.lower()
        self.ollama_url = ollama_url
        self.api_key = api_key or os.getenv('LLM_API_KEY')
        self.selection_strategy = selection_strategy

        # 初始化模型監控器
        self.model_monitor = None
        if MONITOR_AVAILABLE:
            try:
                self.model_monitor = ModelMonitor()
            except Exception as e:
                print(f"⚠️  模型監控器初始化失敗：{e}")
                self.model_monitor = None

        # 載入模型選擇配置
        self.model_config = self._load_model_config()

        # 載入模板
        if template_path is None:
            template_path = Path(__file__).parent.parent.parent / "templates" / "prompts" / "journal_club_template.jinja2"

        with open(template_path, 'r', encoding='utf-8') as f:
            self.template = Template(f.read())

        # 載入風格配置
        if styles_config is None:
            styles_config = Path(__file__).parent.parent.parent / "templates" / "styles" / "academic_styles.yaml"

        # 雲端模型推薦配置
        self.CLOUD_MODEL_RECOMMENDATIONS = {
            'zettelkasten': 'minimax-m2:cloud',  # 230B，深度思考和概念提取
            'research_methods': 'minimax-m2:cloud',  # 需要嚴謹分析
            'literature_review': 'minimax-m2:cloud',  # 需要綜合分析
            'modern_academic': 'gpt-oss:20b-cloud',  # 平衡效能
            'teaching': 'phi3.5:cloud',  # 快速清晰的教學內容
            'default': 'gpt-oss:20b-cloud'  # 通用預設
        }

        with open(styles_config, 'r', encoding='utf-8') as f:
            self.styles_config = yaml.safe_load(f)

        # 初始化LLM客戶端
        self._init_llm_clients()

    def get_style_info(self, style: str) -> Dict[str, Any]:
        """獲取風格信息"""
        styles = self.styles_config.get('styles', {})
        if style not in styles:
            available = ', '.join(styles.keys())
            raise ValueError(f"Unknown style: {style}. Available: {available}")
        return styles[style]

    def get_detail_info(self, detail_level: str) -> Dict[str, Any]:
        """獲取詳細程度信息"""
        details = self.styles_config.get('detail_levels', {})
        if detail_level not in details:
            available = ', '.join(details.keys())
            raise ValueError(f"Unknown detail level: {detail_level}. Available: {available}")
        return details[detail_level]

    def get_language_info(self, language: str) -> Dict[str, Any]:
        """獲取語言信息"""
        languages = self.styles_config.get('languages', {})
        if language not in languages:
            available = ', '.join(languages.keys())
            raise ValueError(f"Unknown language: {language}. Available: {available}")
        return languages[language]

    def generate_prompt(self,
                       topic: str,
                       style: str = "modern_academic",
                       detail_level: str = "standard",
                       language: str = "chinese",
                       slide_count: int = 15,
                       pdf_content: Optional[str] = None,
                       custom_requirements: Optional[str] = None) -> str:
        """
        生成LLM提示詞

        Args:
            topic: 簡報主題
            style: 學術風格
            detail_level: 詳細程度
            language: 語言
            slide_count: 投影片數量
            pdf_content: PDF內容（可選）
            custom_requirements: 自訂要求（可選）

        Returns:
            生成的提示詞
        """
        style_info = self.get_style_info(style)
        detail_info = self.get_detail_info(detail_level)
        lang_info = self.get_language_info(language)

        prompt = self.template.render(
            topic=topic,
            slide_count=slide_count,
            style_name=style_info['name'],
            style_description=style_info['description'],
            detail_name=detail_info['name'],
            detail_description=detail_info['description'],
            detail_level=detail_level,
            language=language,
            pdf_content=pdf_content,
            custom_requirements=custom_requirements
        )

        return prompt

    def _init_llm_clients(self):
        """初始化LLM客戶端（僅初始化指定的provider）"""
        # 只有在 auto 模式或明確指定時才初始化對應的客戶端
        # 這避免了不必要的 API 初始化和錯誤

        # Google Gemini
        if (self.llm_provider in ["auto", "google"]) and GOOGLE_AVAILABLE and (self.api_key or os.getenv('GOOGLE_API_KEY')):
            try:
                genai.configure(api_key=self.api_key or os.getenv('GOOGLE_API_KEY'))
                self.google_client = genai
            except Exception:
                self.google_client = None
        else:
            self.google_client = None

        # OpenAI
        if (self.llm_provider in ["auto", "openai"]) and OPENAI_AVAILABLE and (self.api_key or os.getenv('OPENAI_API_KEY')):
            try:
                self.openai_client = openai.OpenAI(api_key=self.api_key or os.getenv('OPENAI_API_KEY'))
            except Exception:
                self.openai_client = None
        else:
            self.openai_client = None

        # Anthropic Claude
        if (self.llm_provider in ["auto", "anthropic"]) and ANTHROPIC_AVAILABLE and (self.api_key or os.getenv('ANTHROPIC_API_KEY')):
            try:
                self.anthropic_client = anthropic.Anthropic(api_key=self.api_key or os.getenv('ANTHROPIC_API_KEY'))
            except Exception:
                self.anthropic_client = None
        else:
            self.anthropic_client = None

        # OpenRouter: 不需要客戶端初始化，直接使用 requests + OPENROUTER_API_KEY
        # 檢測在 _detect_available_providers() 中通過環境變數完成

    def _load_model_config(self) -> Dict:
        """載入模型選擇配置"""
        config_path = Path(__file__).parent.parent.parent / "config" / "model_selection.yaml"
        if config_path.exists():
            try:
                with open(config_path, 'r', encoding='utf-8') as f:
                    return yaml.safe_load(f)
            except Exception as e:
                print(f"⚠️  載入模型配置失敗：{e}")
        return {}

    def _select_best_model(self, task_type: Optional[str] = None,
                          style: Optional[str] = None) -> Tuple[str, str]:
        """
        智能選擇最佳模型

        Args:
            task_type: 任務類型（zettelkasten/academic_slides等）
            style: 學術風格

        Returns:
            (provider, model_name) 元組
        """
        if not self.model_config:
            # 沒有配置，使用默認
            return self._get_default_model()

        # 獲取可用的提供者
        available_providers = self._detect_available_providers()
        if not available_providers:
            raise RuntimeError("沒有可用的LLM提供者")

        # 根據風格獲取推薦模型
        candidates = []
        if style and "style_model_mapping" in self.model_config:
            style_mapping = self.model_config["style_model_mapping"].get(style, {})
            candidates.extend(style_mapping.get("preferred", []))
            candidates.extend(style_mapping.get("fallback", []))
            avoid = style_mapping.get("avoid", [])
        elif task_type and "task_model_mapping" in self.model_config:
            task_mapping = self.model_config["task_model_mapping"].get(task_type, {})
            candidates.extend(task_mapping.get("preferred", []))
            candidates.extend(task_mapping.get("fallback", []))
            avoid = task_mapping.get("avoid", [])
        else:
            # 使用優先級順序
            models = self.model_config.get("models", {})
            sorted_models = sorted(models.items(),
                                 key=lambda x: x[1].get("priority", 999))
            candidates = [m[0] for m in sorted_models]
            avoid = []

        # 選擇第一個可用的模型
        for model_id in candidates:
            if model_id in avoid:
                continue

            model_info = self.model_config.get("models", {}).get(model_id)
            if not model_info:
                continue

            provider = model_info.get("provider")
            model_name = model_info.get("model_name")

            # 檢查提供者是否可用
            if provider not in available_providers:
                continue

            # 檢查配額（如果有監控器）
            if self.model_monitor:
                quota_status = self.model_monitor.check_quota_status(provider, model_name)
                if quota_status.get("exceeded", False):
                    continue

                # 檢查成本限制
                cost_status = self.model_monitor.check_cost_status()
                if cost_status.get("controlled") and cost_status["session"].get("exceeded"):
                    # 成本超限，選擇免費模型
                    if not model_info.get("free_quota", False):
                        continue

            print(f"🎯 智能選擇模型：{model_id} ({provider}/{model_name})")
            return provider, model_name

        # 沒有找到合適的，使用默認
        return self._get_default_model()

    def _get_default_model(self) -> Tuple[str, str]:
        """獲取默認模型"""
        available = self._detect_available_providers()
        if "google" in available:
            return "google", "gemini-2.0-flash-exp"
        elif "anthropic" in available:
            return "anthropic", "claude-3-haiku-20240307"
        elif "ollama" in available:
            return "ollama", "gpt-oss:20b-cloud"
        else:
            raise RuntimeError("沒有可用的LLM提供者")

    def _check_ollama_health(self) -> bool:
        """檢查Ollama服務健康狀態"""
        try:
            response = requests.get(f"{self.ollama_url}/api/tags", timeout=5)
            return response.status_code == 200
        except Exception:
            return False

    def _detect_available_providers(self) -> List[str]:
        """偵測可用的LLM提供者"""
        providers = []

        # 檢查 Ollama
        if self._check_ollama_health():
            providers.append('ollama')

        # 檢查 Google
        if self.google_client:
            providers.append('google')

        # 檢查 OpenAI
        if self.openai_client:
            providers.append('openai')

        # 檢查 Anthropic
        if self.anthropic_client:
            providers.append('anthropic')

        # 檢查 OpenRouter
        if os.getenv('OPENROUTER_API_KEY'):
            providers.append('openrouter')

        return providers

    def call_llm(self,
                 prompt: str,
                 model: Optional[str] = None,
                 provider: Optional[str] = None,
                 timeout: int = 300,
                 task_type: Optional[str] = None,
                 style: Optional[str] = None,
                 max_tokens: int = 4096) -> Tuple[str, str]:
        """
        統一的LLM調用接口，支援多後端和自動fallback

        Args:
            prompt: 提示詞
            model: 模型名稱
            provider: 指定LLM提供者（可選，None則自動選擇）
            timeout: 超時時間（秒）
            task_type: 任務類型（用於智能選擇）
            style: 學術風格（用於智能選擇）
            max_tokens: 最大生成 tokens 數（僅 OpenRouter，默認 4096）

        Returns:
            (生成的內容, 使用的provider)
        """
        # 決定使用的provider和model
        actual_model = model
        if provider is None:
            if self.llm_provider == "auto":
                # 使用智能選擇
                selected_provider, selected_model = self._select_best_model(task_type, style)
                provider = selected_provider
                if not model:  # 如果沒有指定模型，使用智能選擇的
                    actual_model = selected_model
                print(f"🤖 智能選擇：{provider}/{actual_model}")
            else:
                provider = self.llm_provider

        # 準備fallback chain
        fallback_chain = [provider]
        available = self._detect_available_providers()
        for p in available:
            if p not in fallback_chain:
                fallback_chain.append(p)

        # 嘗試調用LLM（帶fallback）
        last_error = None
        for attempt_provider in fallback_chain:
            start_time = time.time()
            try:
                # 根據提供者調用對應的方法
                if attempt_provider == 'ollama':
                    used_model = actual_model or "gpt-oss:20b-cloud"
                    result = self.call_ollama(prompt, used_model, timeout)
                elif attempt_provider == 'google':
                    used_model = actual_model or "gemini-2.0-flash-exp"
                    result = self.call_google(prompt, used_model)
                elif attempt_provider == 'openai':
                    used_model = actual_model or "gpt-3.5-turbo"
                    result = self.call_openai(prompt, used_model)
                elif attempt_provider == 'anthropic':
                    used_model = actual_model or "claude-3-haiku-20240307"
                    result = self.call_anthropic(prompt, used_model)
                elif attempt_provider == 'openrouter':
                    used_model = actual_model or "anthropic/claude-3.5-sonnet"
                    result = self.call_openrouter(prompt, used_model, timeout, max_tokens)
                else:
                    continue

                # 計算響應時間
                response_time = time.time() - start_time

                # 追蹤使用情況（如果有監控器）
                if self.model_monitor:
                    # 估算token數（簡單估算：字元數/4）
                    estimated_tokens = (len(prompt) + len(result)) // 4
                    usage_info = self.model_monitor.track_usage(
                        model_name=used_model,
                        provider=attempt_provider,
                        tokens_used=estimated_tokens,
                        response_time=response_time,
                        success=True,
                        task_type=task_type
                    )

                    # 檢查是否需要切換模型
                    if usage_info.get("quota_status", {}).get("warning"):
                        print(f"⚠️  配額警告：{used_model} 使用量接近限制")
                    if usage_info.get("cost_status", {}).get("session", {}).get("warning"):
                        print(f"💰 成本警告：當前會話成本 ${usage_info['session_cost']:.2f}")

                return result, attempt_provider

            except Exception as e:
                # 追蹤失敗（如果有監控器）
                if self.model_monitor:
                    response_time = time.time() - start_time
                    self.model_monitor.track_usage(
                        model_name=actual_model or "unknown",
                        provider=attempt_provider,
                        tokens_used=0,
                        response_time=response_time,
                        success=False,
                        task_type=task_type,
                        error_message=str(e)
                    )

                last_error = e
                if len(fallback_chain) > 1:
                    print(f"⚠️  {attempt_provider} 失敗：{str(e)}")
                    if attempt_provider != fallback_chain[-1]:
                        print(f"🔄 嘗試fallback到下一個provider...")
                continue

        # 所有provider都失敗
        raise RuntimeError(f"所有LLM提供者都失敗。最後錯誤：{last_error}")

    def call_ollama(self,
                   prompt: str,
                   model: str = "gpt-oss:20b-cloud",
                   timeout: int = 300) -> str:
        """
        調用Ollama API生成內容（支援本地和雲端模型）

        Args:
            prompt: 提示詞
            model: 模型名稱（支援 :cloud 後綴的雲端模型）
            timeout: 超時時間（秒）

        Returns:
            生成的內容

        支援的雲端模型：
            - minimax-m2:cloud (230B, 支援thinking) - 僅CLI
            - gpt-oss:20b-cloud (20B)
            - qwen2.5-coder:cloud (3.1B, 程式碼分析)
            - phi3.5:cloud (3.8B, 通用任務)
        """
        # 特殊處理 MiniMax-M2:cloud - 使用CLI方式
        if model.lower() == "minimax-m2:cloud":
            import subprocess
            import tempfile
            import os

            try:
                # 將prompt寫入臨時文件以避免命令行長度限制
                with tempfile.NamedTemporaryFile(mode='w', delete=False, suffix='.txt', encoding='utf-8') as f:
                    f.write(prompt)
                    temp_file = f.name

                # 使用ollama CLI運行
                cmd = f'type "{temp_file}" | ollama run {model}'

                result = subprocess.run(
                    cmd,
                    shell=True,
                    capture_output=True,
                    text=True,
                    encoding='utf-8',
                    errors='ignore',  # 忽略編碼錯誤
                    timeout=timeout
                )

                # 清理臨時文件
                try:
                    os.unlink(temp_file)
                except:
                    pass

                if result.returncode != 0:
                    raise RuntimeError(f"Ollama CLI failed: {result.stderr}")

                # 返回輸出，移除thinking部分和ANSI控制字符
                output = result.stdout

                # 移除ANSI轉義序列
                import re
                ansi_escape = re.compile(r'\x1B(?:[@-Z\\-_]|\[[0-?]*[ -/]*[@-~])')
                output = ansi_escape.sub('', output)

                # 提取thinking之後的實際回應
                if "Thinking..." in output:
                    # 找到"...done thinking."之後的內容
                    if "...done thinking." in output:
                        parts = output.split("...done thinking.")
                        if len(parts) > 1:
                            output = parts[1].strip()
                    elif "done thinking" in output:
                        parts = output.split("done thinking")
                        if len(parts) > 1:
                            output = parts[1].strip()

                # 清理殘留的控制字符
                output = re.sub(r'\[.*?\]', '', output)
                output = re.sub(r'[⠙⠹⠸⠴⠦⠧⠏⠋]', '', output)
                output = output.strip()

                return output.strip()

            except subprocess.TimeoutExpired:
                raise RuntimeError(f"Ollama CLI timeout after {timeout} seconds")
            except Exception as e:
                raise RuntimeError(f"Ollama CLI execution failed: {e}")

        # 其他模型使用API方式
        api_endpoint = f"{self.ollama_url}/api/generate"

        payload = {
            "model": model,
            "prompt": prompt,
            "stream": False
        }

        try:
            response = requests.post(
                api_endpoint,
                json=payload,
                timeout=timeout
            )
            response.raise_for_status()

            result = response.json()
            return result.get('response', '')

        except requests.exceptions.RequestException as e:
            raise RuntimeError(f"Ollama API call failed: {e}")

    def call_google(self,
                   prompt: str,
                   model: str = "gemini-pro") -> str:
        """
        調用Google Gemini API生成內容

        Args:
            prompt: 提示詞
            model: 模型名稱

        Returns:
            生成的內容
        """
        if not self.google_client:
            raise RuntimeError("Google Gemini not available. Install: pip install google-generativeai")

        try:
            model_instance = self.google_client.GenerativeModel(model)
            response = model_instance.generate_content(prompt)
            return response.text
        except Exception as e:
            raise RuntimeError(f"Google Gemini API call failed: {e}")

    def call_openai(self,
                   prompt: str,
                   model: str = "gpt-3.5-turbo") -> str:
        """
        調用OpenAI API生成內容

        Args:
            prompt: 提示詞
            model: 模型名稱

        Returns:
            生成的內容
        """
        if not self.openai_client:
            raise RuntimeError("OpenAI not available. Install: pip install openai")

        try:
            response = self.openai_client.chat.completions.create(
                model=model,
                messages=[
                    {"role": "user", "content": prompt}
                ]
            )
            return response.choices[0].message.content
        except Exception as e:
            raise RuntimeError(f"OpenAI API call failed: {e}")

    def call_anthropic(self,
                      prompt: str,
                      model: str = "claude-3-sonnet-20240229") -> str:
        """
        調用Anthropic Claude API生成內容

        Args:
            prompt: 提示詞
            model: 模型名稱

        Returns:
            生成的內容
        """
        if not self.anthropic_client:
            raise RuntimeError("Anthropic Claude not available. Install: pip install anthropic")

        try:
            response = self.anthropic_client.messages.create(
                model=model,
                max_tokens=4096,
                messages=[
                    {"role": "user", "content": prompt}
                ]
            )
            return response.content[0].text
        except Exception as e:
            raise RuntimeError(f"Anthropic Claude API call failed: {e}")

    def call_openrouter(self,
                       prompt: str,
                       model: str = "anthropic/claude-3.5-sonnet",
                       timeout: int = 300,
                       max_tokens: int = 4096) -> str:
        """
        調用 OpenRouter API 生成內容

        Args:
            prompt: 提示詞
            model: 模型名稱（例如: anthropic/claude-3.5-sonnet）
            timeout: 超時時間（秒）
            max_tokens: 最大生成 tokens 數（默認 4096）

        Returns:
            生成的內容

        支援的模型範例：
            - anthropic/claude-3.5-sonnet (推薦用於 Zettelkasten)
            - anthropic/claude-3-haiku (快速經濟)
            - google/gemini-2.0-flash-exp (免費)
        """
        api_key = os.getenv('OPENROUTER_API_KEY')
        if not api_key:
            raise ValueError("OPENROUTER_API_KEY not set. Please add it to .env file")

        url = "https://openrouter.ai/api/v1/chat/completions"

        headers = {
            "Authorization": f"Bearer {api_key}",
            "HTTP-Referer": "https://github.com/claude-lit-workflow",
            "X-Title": "Claude Lit Workflow - Zettelkasten Generator"
        }

        data = {
            "model": model,
            "messages": [
                {"role": "user", "content": prompt}
            ],
            "max_tokens": max_tokens
        }

        try:
            response = requests.post(url, headers=headers, json=data, timeout=timeout)
            response.raise_for_status()

            result = response.json()
            return result['choices'][0]['message']['content']

        except requests.exceptions.Timeout:
            raise RuntimeError(f"OpenRouter API call timeout after {timeout}s")
        except requests.exceptions.RequestException as e:
            raise RuntimeError(f"OpenRouter API call failed: {e}")
        except (KeyError, IndexError) as e:
            raise RuntimeError(f"Failed to parse OpenRouter response: {e}")

    def parse_slides(self, content: str) -> List[Dict[str, str]]:
        """
        解析LLM生成的投影片內容

        Args:
            content: LLM生成的文本

        Returns:
            投影片列表，每個投影片包含 title 和 content
        """
        # 使用正則表達式分割投影片
        pattern = r'===([^=]+)==='
        parts = re.split(pattern, content)

        slides = []

        # 跳過第一部分（通常是前置文本）
        for i in range(1, len(parts), 2):
            if i + 1 < len(parts):
                title = parts[i].strip()
                slide_content = parts[i + 1].strip()

                slides.append({
                    'title': title,
                    'content': slide_content,
                    'type': 'title' if '標題頁' in title or 'Title' in title else 'content'
                })

        return slides

    def create_markdown(self,
                       slides: List[Dict[str, str]],
                       output_path: str,
                       title: Optional[str] = None,
                       style: str = "modern_academic") -> str:
        """
        創建Markdown簡報文件（支援Marp/reveal.js格式）

        Args:
            slides: 投影片數據
            output_path: 輸出路徑
            title: 簡報標題
            style: 學術風格

        Returns:
            輸出文件路徑
        """
        # 載入Markdown簡報模板
        md_template_path = Path(__file__).parent.parent.parent / "templates" / "markdown" / "academic_slides.jinja2"

        with open(md_template_path, 'r', encoding='utf-8') as f:
            md_template = Template(f.read())

        # 準備數據
        main_title = title or "學術簡報"
        subtitle = ""
        content_slides = []

        for slide_data in slides:
            if slide_data['type'] == 'title':
                # 解析標題頁
                lines = slide_data['content'].split('\n')
                for line in lines:
                    line = line.strip()
                    if line.startswith('標題：') or line.startswith('Title:'):
                        main_title = line.split('：', 1)[-1].split(':', 1)[-1].strip()
                    elif line.startswith('副標題：') or line.startswith('Subtitle:'):
                        subtitle = line.split('：', 1)[-1].split(':', 1)[-1].strip()
            else:
                # 解析內容頁
                content_lines = slide_data['content'].split('\n')
                slide_title = slide_data['title']
                items = []

                for line in content_lines:
                    line = line.strip()
                    if not line:
                        continue
                    if line.startswith('標題：') or line.startswith('Title:'):
                        slide_title = line.split('：', 1)[-1].split(':', 1)[-1].strip()
                    elif line.startswith('內容：') or line.startswith('Content:'):
                        continue
                    else:
                        # 清理項目符號
                        line = re.sub(r'^[•\-\*]\s*', '', line)
                        if line:
                            items.append(line)

                content_slides.append({
                    'title': slide_title,
                    'items': items,
                    'notes': ''
                })

        # 渲染模板
        current_date = datetime.now().strftime("%Y-%m-%d")

        markdown_content = md_template.render(
            main_title=main_title,
            subtitle=subtitle,
            slides=content_slides,
            header_text=style,
            footer_text=current_date,
            authors='',
            date=current_date,
            contact=''
        )

        # 保存文件
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)

        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(markdown_content)

        return str(output_path)

    def create_pptx(self,
                   slides: List[Dict[str, str]],
                   output_path: str,
                   title: Optional[str] = None) -> str:
        """
        創建PowerPoint文件

        Args:
            slides: 投影片數據
            output_path: 輸出路徑
            title: 簡報標題（可選）

        Returns:
            輸出文件路徑
        """
        prs = Presentation()

        # 設置投影片大小（16:9）
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)

        for slide_data in slides:
            if slide_data['type'] == 'title':
                # 標題頁
                slide = prs.slides.add_slide(prs.slide_layouts[0])

                # 解析標題和副標題
                lines = slide_data['content'].split('\n')
                main_title = ""
                subtitle = ""

                for line in lines:
                    line = line.strip()
                    if line.startswith('標題：') or line.startswith('Title:'):
                        main_title = line.split('：', 1)[-1].split(':', 1)[-1].strip()
                    elif line.startswith('副標題：') or line.startswith('Subtitle:'):
                        subtitle = line.split('：', 1)[-1].split(':', 1)[-1].strip()

                if slide.shapes.title:
                    slide.shapes.title.text = main_title or slide_data['title']

                if len(slide.placeholders) > 1:
                    slide.placeholders[1].text = subtitle

            else:
                # 內容頁
                slide = prs.slides.add_slide(prs.slide_layouts[1])

                # 解析內容以提取真正的標題
                content_lines = slide_data['content'].split('\n')
                actual_title = slide_data['title']  # 預設值
                content_items = []

                for line in content_lines:
                    line = line.strip()
                    if not line:
                        continue

                    # 提取標題
                    if line.startswith('標題：') or line.startswith('Title:'):
                        actual_title = line.split('：', 1)[-1].split(':', 1)[-1].strip()
                    # 跳過 "內容：" 標記行
                    elif line.startswith('內容：') or line.startswith('Content:'):
                        continue
                    # 收集實際內容
                    else:
                        content_items.append(line)

                # 設置標題
                if slide.shapes.title:
                    slide.shapes.title.text = actual_title

                # 添加內容
                if len(slide.placeholders) > 1:
                    text_frame = slide.placeholders[1].text_frame
                    text_frame.clear()

                    # 啟用文字自動調整和換行
                    text_frame.word_wrap = True

                    # 計算總內容長度以決定字體大小
                    total_content_length = sum(len(item) for item in content_items)
                    item_count = len(content_items)

                    # 智能字體大小決策
                    if total_content_length > 1000 or item_count > 8:
                        base_font_size = Pt(11)
                        line_spacing = 0.9
                    elif total_content_length > 800 or item_count > 6:
                        base_font_size = Pt(12)
                        line_spacing = 1.0
                    elif total_content_length > 600 or item_count > 5:
                        base_font_size = Pt(14)
                        line_spacing = 1.1
                    elif total_content_length > 400:
                        base_font_size = Pt(16)
                        line_spacing = 1.2
                    else:
                        base_font_size = Pt(18)
                        line_spacing = 1.3

                    for line in content_items:
                        # 移除項目符號標記
                        line = re.sub(r'^[•\-\*]\s*', '', line)
                        # 移除粗體標記 **text**
                        line = line.replace('**', '')

                        if line:
                            p = text_frame.add_paragraph()
                            p.text = line
                            p.level = 0
                            p.font.size = base_font_size
                            p.space_before = Pt(3)
                            p.space_after = Pt(3)
                            p.line_spacing = line_spacing

        # 保存文件
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))

        return str(output_path)

    def generate_slides(self,
                       topic: str,
                       style: str = "modern_academic",
                       detail_level: str = "standard",
                       language: str = "chinese",
                       slide_count: int = 15,
                       output_path: Optional[str] = None,
                       output_format: str = "pptx",
                       pdf_content: Optional[str] = None,
                       custom_requirements: Optional[str] = None,
                       model: str = "gemma2:latest") -> Dict[str, Any]:
        """
        完整的投影片生成流程

        Args:
            topic: 簡報主題
            style: 學術風格
            detail_level: 詳細程度
            language: 語言
            slide_count: 投影片數量
            output_path: 輸出路徑
            output_format: 輸出格式 (pptx/markdown/both)
            pdf_content: PDF內容
            custom_requirements: 自訂要求
            model: Ollama模型

        Returns:
            包含結果信息的字典
        """
        # 1. 生成提示詞
        prompt = self.generate_prompt(
            topic=topic,
            style=style,
            detail_level=detail_level,
            language=language,
            slide_count=slide_count,
            pdf_content=pdf_content,
            custom_requirements=custom_requirements
        )

        # 2. 調用LLM（自動選擇可用的provider）
        print("🤖 正在生成投影片內容...")
        llm_output, used_provider = self.call_llm(prompt, model=model)
        print(f"✅ 使用 {used_provider} 生成完成")

        # 3. 解析投影片
        print("📊 正在解析投影片結構...")
        slides = self.parse_slides(llm_output)

        if not slides:
            raise ValueError("無法解析投影片內容，請檢查LLM輸出格式")

        # 4. 生成輸出文件
        if output_path is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            base_name = f"output/{topic}_{style}_{timestamp}"
        else:
            base_name = str(Path(output_path).with_suffix(''))

        output_files = []

        # 生成PPTX
        if output_format in ["pptx", "both"]:
            pptx_path = f"{base_name}.pptx"
            print("💾 正在創建PowerPoint文件...")
            pptx_file = self.create_pptx(slides, pptx_path, title=topic)
            output_files.append(pptx_file)

        # 生成Markdown
        if output_format in ["markdown", "both"]:
            md_path = f"{base_name}.md"
            print("📝 正在創建Markdown簡報...")
            md_file = self.create_markdown(slides, md_path, title=topic, style=style)
            output_files.append(md_file)

        return {
            'success': True,
            'output_path': output_files[0] if len(output_files) == 1 else output_files,
            'output_files': output_files,
            'output_format': output_format,
            'slide_count': len(slides),
            'style': style,
            'detail_level': detail_level,
            'language': language,
            'llm_provider': used_provider,
            'llm_output': llm_output[:500]  # 保存前500字元作為預覽
        }


# 便捷函數
def make_slides(topic: str,
               pdf_path: Optional[str] = None,
               style: str = "modern_academic",
               detail_level: str = "standard",
               language: str = "chinese",
               slide_count: int = 15,
               output_path: Optional[str] = None,
               model: str = "gemma2:latest") -> str:
    """
    便捷函數：生成投影片

    Args:
        topic: 簡報主題
        pdf_path: PDF文件路徑（可選）
        style: 學術風格
        detail_level: 詳細程度
        language: 語言
        slide_count: 投影片數量
        output_path: 輸出路徑
        model: Ollama模型

    Returns:
        輸出PPTX文件路徑
    """
    maker = SlideMaker()

    # 如果提供PDF路徑，提取內容
    pdf_content = None
    if pdf_path:
        from ..extractors import PDFExtractor
        extractor = PDFExtractor(max_chars=10000)  # Journal Club限制
        result = extractor.extract(pdf_path)
        pdf_content = result['full_text']

    result = maker.generate_slides(
        topic=topic,
        style=style,
        detail_level=detail_level,
        language=language,
        slide_count=slide_count,
        output_path=output_path,
        pdf_content=pdf_content,
        model=model
    )

    return result['output_path']


if __name__ == "__main__":
    # 測試代碼
    import sys

    if len(sys.argv) > 1:
        topic = sys.argv[1]
        print(f"測試生成投影片: {topic}")

        maker = SlideMaker()
        result = maker.generate_slides(
            topic=topic,
            style="modern_academic",
            slide_count=5  # 測試用少量投影片
        )

        print(f"\n✅ 完成！")
        print(f"輸出: {result['output_path']}")
        print(f"投影片數: {result['slide_count']}")
    else:
        print("用法: python slide_maker.py <主題>")
